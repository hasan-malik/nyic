#!/usr/bin/env python3
"""Generate ichooseny-deck.pptx — NYIC #IChooseNY pitch deck.

Follows the official NYIC "Pitch Day Starter Deck" template structure,
kept deliberately sparse: visual-first, <6 min, "show fewer things but make
them feel real". Instructional Tip boxes are intentionally omitted; real
product screenshots from screenshots/ are dropped into the image slots.

Requires python-pptx. Run: python3 scripts/gen_pptx.py
"""
import os
import struct
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- NYIC brand palette --------------------------------------------
PURPLE = RGBColor(0x8A, 0x0F, 0xE6)   # signature violet (backgrounds)
BLACK  = RGBColor(0x00, 0x00, 0x00)
HEADP  = RGBColor(0xB0, 0x47, 0xFF)   # section labels on black
ORANGE = RGBColor(0xF2, 0x6B, 0x21)   # sub-headers
YELLOW = RGBColor(0xE8, 0xFB, 0x3C)   # accent labels
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0xBF, 0xBF, 0xBF)   # placeholders
DGRAY  = RGBColor(0x3A, 0x3A, 0x3A)   # footer glyph
FRAME  = RGBColor(0x2A, 0x2A, 0x2A)   # image frame on black

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "screenshots")
OUT = os.path.join(ROOT, "ichooseny-deck.pptx")

# Map semantic keys -> a unique timestamp substring of the real screenshot
# filename (the files use a narrow no-break space, so we match loosely).
IMG = {
    "stories":       "8.53.29",  # "3,761 voices and counting" story bank
    "dashboard":     "8.54.17",  # staff Overview console + AI synthesis
    "origins":       "7.13.07",  # world "Every origin. One city." map
    "constellation": "8.53.35",  # "Every voice. Every thread." kinship galaxy
}


def resolve(key):
    """Return the absolute path of the screenshot matching key's timestamp."""
    needle = IMG.get(key, key)
    for f in os.listdir(SHOTS):
        if needle in f:
            return os.path.join(SHOTS, f)
    raise FileNotFoundError("no screenshot matching %r" % needle)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---- helpers -------------------------------------------------------
def slide(bg=BLACK):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def P(*runs, align=PP_ALIGN.LEFT, space_after=4):
    return {"runs": list(runs), "align": align, "space_after": space_after}


def text(s, l, t, w, h, paras, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para["align"]
        p.space_after = Pt(para["space_after"])
        for run in para["runs"]:
            r = p.add_run()
            r.text = run[0]
            f = r.font
            f.size = Pt(run[1]); f.bold = run[2]; f.color.rgb = run[3]
            f.italic = run[4] if len(run) > 4 else False
            f.name = "Arial"
    return tb


def rect(s, l, t, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(24)
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def picture(s, key, l, t, w, h, frame=FRAME):
    """Contain a screenshot inside box (l,t,w,h) inches, centered, framed."""
    path = resolve(key)
    iw, ih = png_size(path)
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = l + (w - pw) / 2, t + (h - ph) / 2
    if frame is not None:
        fr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(px - 0.03), Inches(py - 0.03),
                                Inches(pw + 0.06), Inches(ph + 0.06))
        fr.fill.solid(); fr.fill.fore_color.rgb = frame
        fr.line.fill.background(); fr.shadow.inherit = False
    s.shapes.add_picture(path, Inches(px), Inches(py), Inches(pw), Inches(ph))


def footer(s):
    text(s, 0.45, 6.82, 4, 0.5,
         [P(("➜   ●   #   ", 16, True, DGRAY),
            ("n", 18, True, PURPLE),
            ("   ◼", 16, True, DGRAY))])


def bullets(s, l, t, w, items, size=18, color=WHITE, gap=10):
    paras = [P(("•  " + it, size, False, color), space_after=gap) for it in items]
    text(s, l, t, w, 4, paras)


def placeholder(s, l, t, w, h, label):
    rect(s, l, t, w, h, GRAY)
    text(s, l, t + h / 2 - 0.35, w, 0.7,
         [P((label, 22, True, BLACK), align=PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)


# ===================================================================
# SLIDE 1 — Cover page
# ===================================================================
s = slide(PURPLE)
text(s, 0.7, 1.25, 9.8, 2.6, [
    P(("#IChooseNY", 60, True, WHITE, True), space_after=10),
    P(("A living story bank that turns 40 years of immigrant "
       "stories into action.", 30, True, WHITE, True)),
])
text(s, 8.2, 2.7, 4.4, 0.5, [P(("Core Idea", 18, True, YELLOW), align=PP_ALIGN.RIGHT)])
text(s, 0.7, 5.45, 7, 0.5, [P(("New York Immigration Coalition", 22, False, WHITE))])
text(s, 8.2, 5.45, 4.4, 0.5, [P(("Nonprofit Partner", 18, True, YELLOW), align=PP_ALIGN.RIGHT)])
text(s, 0.7, 6.05, 7, 0.5, [P(("Team [Your Team Name]", 22, False, WHITE))])
text(s, 8.2, 6.05, 4.4, 0.5, [P(("Team Name", 18, True, YELLOW), align=PP_ALIGN.RIGHT)])
footer(s)

# ===================================================================
# SLIDE 2 — The Opportunity / Vision
# ===================================================================
s = slide(BLACK)
text(s, 0.7, 0.7, 11, 0.6, [P(("THE OPPORTUNITY / VISION", 18, True, HEADP))])
text(s, 0.9, 2.3, 11.5, 3, [
    P(("Every immigrant story, heard — and turned into power.", 40, True, WHITE, True),
      space_after=14),
    P(("NYIC's 40 years become a living movement, not an archive.", 26, False, GRAY)),
])
footer(s)

# ===================================================================
# SLIDE 3 — System Overview (anchor)
# ===================================================================
s = slide(BLACK)
text(s, 0.7, 0.55, 11.8, 0.6, [P(("SYSTEM OVERVIEW", 18, True, HEADP))])
text(s, 0.7, 1.15, 11.8, 0.6,
     [P(("Five minutes of audio → a living archive → daily advocacy.", 20, False, WHITE))])

cols = [
    ("Story Collection\n& Campaign", "Entry & Awareness",
     ["Friend-to-friend audio", "Phone, QR, or staff-led", "Consent built in"]),
    ("Story Bank\n& Internal System", "Internal Management",
     ["Auto-transcribe + tag", "Searchable, 8 dimensions", "Care-first review"]),
    ("Story Activation,\nEngagement & Action", "Visualization & CTA",
     ["One story a day", "Match stories to bills", "Donate · volunteer · testify"]),
]
cw, gap, x0 = 3.7, 0.45, 0.7
for i, (title, sub, items) in enumerate(cols):
    x = x0 + i * (cw + gap)
    chev = rect(s, x, 2.15, cw, 0.95, PURPLE, shape=MSO_SHAPE.CHEVRON)
    tf = chev.text_frame; tf.word_wrap = True
    p0 = tf.paragraphs[0]
    for j, line in enumerate(title.split("\n")):
        pp = p0 if j == 0 else tf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        r = pp.add_run(); r.text = line
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Arial"
    text(s, x, 3.35, cw, 0.5, [P((sub, 16, True, ORANGE))])
    bullets(s, x, 3.95, cw, items, size=15, gap=8)
footer(s)

# ===================================================================
# SLIDES 4-6 — Three stages (text left, real screenshot right)
# ===================================================================
stages = [
    ("STORY COLLECTION EXPERIENCE",
     ["A friend interviews a friend", "5-min audio, any language",
      "QR codes, events, social", "Consent before anything"],
     "stories"),
    ("INTERNAL STORY BANK / AI SYSTEM",
     ["Auto-transcribe + translate", "Tagged across 8 dimensions",
      "Search & approve in one console", "~2¢ per story"],
     "dashboard"),
    ("STORY ACTIVATION & COMMUNITY ACTION",
     ["One story a day → Instagram", "“A bill is moving” → matched stories",
      "Map of all of New York", "Donate · volunteer · testify"],
     "origins"),
]
for head, items, shot in stages:
    s = slide(BLACK)
    text(s, 0.7, 0.7, 6, 0.6, [P((head, 18, True, HEADP))])
    text(s, 0.7, 1.6, 5.5, 0.5, [P(("Focus", 16, True, ORANGE))])
    bullets(s, 0.7, 2.2, 5.5, items, size=20, gap=16)
    picture(s, shot, 6.55, 1.5, 6.3, 4.6)
    footer(s)

# ===================================================================
# SLIDE 7 — 15-Sec Concept Video
# ===================================================================
s = slide(BLACK)
placeholder(s, 3.4, 1.1, 6.5, 5.3, "▶  15-Sec Concept Video")
footer(s)

# ===================================================================
# SLIDE 8 — Impact & Closing
# ===================================================================
s = slide(BLACK)
text(s, 0.7, 0.55, 11.8, 0.6, [P(("IMPACT & CLOSING", 18, True, HEADP))])
text(s, 0.9, 1.35, 11.5, 1.8, [
    P(("“We don't ask immigrants to prove they belong.", 30, True, WHITE, True), space_after=4),
    P(("We ask them to choose us back.”", 30, True, WHITE, True)),
])
text(s, 0.9, 3.15, 11.5, 0.7,
     [P(("3,700+ stories", 22, True, YELLOW),
        ("    ·    120+ countries    ·    20 languages    ·    ≈ $79 to process the archive",
         22, False, WHITE))])
picture(s, "constellation", 0.9, 3.95, 11.5, 3.0)
footer(s)

# ===================================================================
# SLIDE 9 — Thank you
# ===================================================================
s = slide(PURPLE)
text(s, 0.7, 1.4, 9, 1.6, [P(("Thank you", 60, True, WHITE, True))])
text(s, 0.7, 5.6, 7.5, 0.6, [P(("#IChooseNY  ·  ichooseny.netlify.app", 22, False, WHITE))])
text(s, 8.2, 5.6, 4.4, 0.6, [P(("Team [Your Team Name]", 20, True, YELLOW), align=PP_ALIGN.RIGHT)])
footer(s)

prs.save(OUT)
print("wrote", OUT, "—", len(prs.slides._sldIdLst), "slides")
