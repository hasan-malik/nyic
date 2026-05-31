#!/usr/bin/env python3
"""Generate ichooseny-deck.pptx — a branded, editable deck for the #IChooseNY pitch."""
import struct
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- brand ----
NAVY = RGBColor(0x0D, 0x1B, 0x4C)
INK = RGBColor(0x0A, 0x13, 0x30)
RED = RGBColor(0xE0, 0x49, 0x2F)
GOLD = RGBColor(0xF0, 0xB4, 0x29)
EMER = RGBColor(0x34, 0xD3, 0x99)
MIST = RGBColor(0xF4, 0xF6, 0xFB)
LBLUE = RGBColor(0xC3, 0xCC, 0xEC)
PEACH = RGBColor(0xFF, 0xD9, 0xCF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SLATE = RGBColor(0x33, 0x41, 0x55)
GREY = RGBColor(0x94, 0xA3, 0xB8)
HEAD = "Montserrat"
BODY = "Inter"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def png_size(path):
    with open(path, "rb") as f:
        head = f.read(24)
    w, h = struct.unpack(">II", head[16:24])
    return w, h


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE, line_w=1.0):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=1.0):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        for (txt, size, color, bold, font) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def picture(s, path, x, y, w, h):
    """Contain image within box (x,y,w,h) inches, centered, with a subtle frame."""
    iw, ih = png_size(path)
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    rect(s, px - 0.04, py - 0.04, pw + 0.08, ph + 0.08, fill=None, line=RGBColor(0x2A, 0x3A, 0x7A), line_w=1.0)
    s.shapes.add_picture(path, Inches(px), Inches(py), Inches(pw), Inches(ph))


def H(txt, size, color=WHITE):
    return (txt, size, color, True, HEAD)


def B(txt, size, color=SLATE, bold=False):
    return (txt, size, color, bold, BODY)


def header_bar(s, label, sub=None, dark=False):
    rect(s, 0, 0, SW, 1.15, fill=INK if dark else NAVY)
    rect(s, 0.6, 0.42, 0.16, 0.34, fill=GOLD)
    runs = [[H(label, 26, WHITE)]]
    text(s, 0.95, 0.3, 9, 0.7, runs, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(s, 0.95, 0.74, 11, 0.35, [[B(sub, 11, RGBColor(0xC3, 0xCC, 0xEC))]])


IMG = "/tmp/deck-{}.png"

# ============ SLIDE 1 — TITLE ============
s = slide(NAVY)
picture(s, IMG.format("origins"), 7.0, 1.7, 6.1, 4.4)
rect(s, 6.6, 0, SW - 6.6, SH, fill=NAVY)  # dim the right image edge into the bg
picture(s, IMG.format("origins"), 7.15, 2.0, 5.9, 3.9)
rect(s, 0.9, 2.0, 0.18, 1.7, fill=GOLD)
text(s, 1.25, 1.9, 6.2, 3.4, [
    [H("#IChooseNY", 54, GOLD)],
    [H("#IChooseYou", 30, WHITE)],
    [B("Belonging isn't a transaction.", 20, RGBColor(0xC3, 0xCC, 0xEC), True)],
], space=1.05)
text(s, 1.25, 6.4, 8, 0.5, [[B("NYIC Social Impact Hackathon  ·  [Your team name]", 13, GREY)]])
notes(s, "Before we say anything — listen. Go STRAIGHT into the 30-second real interview audio. "
         "Do not read the slide. Let it breathe, then begin.")

# ============ SLIDE 2 — TEAM ============
s = slide(MIST)
header_bar(s, "Who we are", "We listened to NYIC's narrative — then built its next chapter.")
roles = [("[Name]", "Product"), ("[Name]", "Design"), ("[Name]", "Data / AI"), ("[Name]", "Storytelling")]
cw = 2.7
x0 = (SW - (cw * 4 + 0.4 * 3)) / 2
for i, (nm, role) in enumerate(roles):
    x = x0 + i * (cw + 0.4)
    rect(s, x, 2.2, cw, 2.7, fill=WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    circ = rect(s, x + cw / 2 - 0.6, 2.55, 1.2, 1.2, fill=NAVY, shape=MSO_SHAPE.OVAL)
    text(s, x, 3.95, cw, 0.5, [[H(nm, 18, NAVY)]], align=PP_ALIGN.CENTER)
    text(s, x, 4.4, cw, 0.4, [[B(role, 13, RED, True)]], align=PP_ALIGN.CENTER)
notes(s, "20 seconds. Names + the one promise: we didn't replace NYIC's message, we evolved it. Move fast.")

# ============ SLIDE 3 — CONCEPT ============
s = slide(NAVY)
text(s, 1.0, 0.8, 11.3, 1.6, [[H("The question a system asks is a statement about what it thinks a person is ", 26, WHITE),
                               ("for", 26, GOLD, True, HEAD), (".", 26, WHITE, True, HEAD)]], space=1.05)
rect(s, 1.0, 2.9, 5.4, 3.4, fill=INK, line=RGBColor(0x2A, 0x3A, 0x7A))
text(s, 1.35, 3.2, 4.8, 2.8, [
    [B("Everyone else asks:", 15, GREY, True)],
    [H("“What did you contribute?”", 24, WHITE)],
    [B("Usefulness only buys tolerance. You may stay as long as you produce.", 14, RGBColor(0xAE, 0xB9, 0xDC))],
], space=1.05)
rect(s, 6.9, 2.9, 5.4, 3.4, fill=RED)
text(s, 7.25, 3.2, 4.8, 2.8, [
    [B("We ask:", 15, PEACH, True)],
    [H("“What did you hold onto that you were told to give up?”", 23, WHITE)],
    [B("#IChooseYou — un-extractable by design.", 14, WHITE, True)],
], space=1.05)
notes(s, "The narrative shift. Usefulness only ever buys tolerance — you stay as long as you produce. "
         "We built a system where belonging is the premise, not the prize. This EVOLVES NYIC's 'we built New York'.")

# ============ SLIDE 4 — FILM ============
s = slide(INK)
rect(s, 3.5, 2.35, 6.3, 2.8, fill=NAVY, line=GOLD, line_w=1.5)
tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(6.35), Inches(3.35), Inches(0.7), Inches(0.8))
tri.rotation = 90
tri.fill.solid(); tri.fill.fore_color.rgb = GOLD; tri.line.fill.background(); tri.shadow.inherit = False
text(s, 3.5, 4.2, 6.3, 0.8, [[H("15-second concept film", 20, WHITE)]], align=PP_ALIGN.CENTER)
text(s, 0, 6.4, SW, 0.5, [[B("Insert > Video here in Google Slides / Canva  ·  file: ichooseny-teaser.mp4 (see HERO_VIDEO.md)", 12, GREY)]], align=PP_ALIGN.CENTER)
notes(s, "Play the 15-sec teaser. Say nothing during it. After: 'That's the feeling. Here's the machine that makes it scale.'")

# ============ SLIDE 5 — JOURNEY ============
s = slide(MIST)
header_bar(s, "Journey of a story", "Five minutes of audio → a living archive")
steps = [("Record", "friend → friend"), ("Transcribe", "ElevenLabs"), ("Care-check", "Claude Haiku"),
         ("Tag", "Claude Sonnet"), ("Connect", "embeddings"), ("Activate", "IG · Runway")]
bw, gap = 1.78, 0.26
x0 = (SW - (bw * 6 + gap * 5)) / 2
for i, (t, sub) in enumerate(steps):
    x = x0 + i * (bw + gap)
    col = [NAVY, RED, GOLD, NAVY, EMER, RED][i]
    rect(s, x, 2.9, bw, 1.5, fill=WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    rect(s, x, 2.9, bw, 0.12, fill=col)
    text(s, x, 3.15, bw, 0.6, [[H(t, 15, NAVY)]], align=PP_ALIGN.CENTER)
    text(s, x, 3.75, bw, 0.5, [[B(sub, 10.5, GREY)]], align=PP_ALIGN.CENTER)
    if i < 5:
        text(s, x + bw - 0.04, 3.2, 0.34, 0.5, [[H("›", 20, GREY)]], align=PP_ALIGN.CENTER)
text(s, 0, 5.1, SW, 1.4, [
    [H("One story. Two cents. Fully automatic after the interview.", 22, NAVY)],
    [B("A daughter interviews her mother for five minutes — everything after is the pipeline.", 14, SLATE)],
], align=PP_ALIGN.CENTER, space=1.1)
notes(s, "Walk the path of ONE story (Amina). Everything after the 5-min interview is automatic — and costs two cents.")

# ============ SLIDE 6 — PUBLIC DEMO ============
s = slide(INK)
header_bar(s, "A bank you can feel, not scroll", "Threads (kinship) · Boroughs & State-wide · 120+ origins → one city", dark=True)
picture(s, IMG.format("origins"), 0.45, 1.35, 6.3, 5.7)
picture(s, IMG.format("statewide"), 6.95, 1.35, 3.05, 2.75)
picture(s, IMG.format("threads"), 6.95, 4.25, 3.05, 2.75)
text(s, 10.2, 1.35, 2.9, 5.7, [
    [H("Real maps.", 18, GOLD)],
    [B("Dark CARTO / OpenStreetMap — open-source, the same stack pros use.", 12.5, LBLUE)],
    [B("", 6, WHITE)],
    [H("All of NY.", 18, GOLD)],
    [B("Buffalo, Rochester, Long Island, the Hudson Valley — not just five boroughs.", 12.5, LBLUE)],
], anchor=MSO_ANCHOR.MIDDLE, space=1.08)
notes(s, "NYIC serves all of New York — so do we. Click a glowing point, play a voice. Show Origins arcs converging on NYC.")

# ============ SLIDE 7 — STAFF DEMO ============
s = slide(INK)
header_bar(s, "The team's daily cockpit", "Care-first review · live analytics · activation · cost transparency", dark=True)
picture(s, IMG.format("dash"), 0.45, 1.35, 8.7, 5.75)
text(s, 9.45, 1.5, 3.5, 5.5, [
    [H("8 tag dimensions", 17, GOLD)],
    [B("mood · theme · topic · struggle · cultural heritage · neighborhood · member org · life stage", 12.5, LBLUE)],
    [B("", 6, WHITE)],
    [H("A live needs-assessment", 17, GOLD)],
    [B("“Struggles named” + “Cultural heritage” charts turn stories into advocacy data.", 12.5, LBLUE)],
], anchor=MSO_ANCHOR.MIDDLE, space=1.08)
notes(s, "Every story tagged across 8 dimensions, so the dashboard doubles as a live needs-assessment for advocacy & services.")

# ============ SLIDE 8 — CARE & CONSENT ============
s = slide(MIST)
header_bar(s, "People over data. Always.")
items = [
    ("Granular consent", "Name, pseudonym, or fully anonymous — chosen per story."),
    ("Voice protection", "ElevenLabs voice-masking for narrators at risk."),
    ("Care, not censorship", "Claude Haiku flags trauma to a human — never auto-rejects."),
    ("Credit where due", "Member-org attribution + role-based access for 200+ orgs."),
]
for i, (t, d) in enumerate(items):
    y = 1.7 + i * 1.28
    rect(s, 1.4, y, 0.16, 1.0, fill=[NAVY, RED, GOLD, EMER][i])
    text(s, 1.8, y, 10.2, 1.1, [
        [H(t, 18, NAVY)],
        [B(d, 14, SLATE)],
    ], space=1.05)
picture(s, IMG.format("detail"), 8.4, 1.7, 4.5, 5.1)
notes(s, "Two doors: tell your story as you, or under protection. We never trade safety for a testimonial. "
         "Directly answers NYIC's StoryCorps / UnLocal data-security concerns.")

# ============ SLIDE 9 — ACTIVATION ============
s = slide(NAVY)
text(s, 1.0, 0.9, 11, 1.0, [[H("One story a day. A movement a year.", 30, WHITE)]])
acts = [
    ("Daily feature", "The system surfaces one story a day, ready to post to Instagram + LinkedIn."),
    ("Runway animation", "Anonymized voices become animated shorts — safe and shareable."),
    ("Advocacy match", "“A bill is moving” → instantly find consenting stories that fit."),
    ("Viral loop", "Every narrator nominates two friends. The bank grows itself."),
]
for i, (t, d) in enumerate(acts):
    x = 1.0 + (i % 2) * 5.9
    y = 2.4 + (i // 2) * 2.1
    rect(s, x, y, 5.4, 1.8, fill=INK, line=RGBColor(0x2A, 0x3A, 0x7A))
    rect(s, x, y, 0.14, 1.8, fill=GOLD)
    text(s, x + 0.4, y + 0.25, 4.8, 1.4, [
        [H(t, 17, GOLD)],
        [B(d, 13.5, LBLUE)],
    ], space=1.05)
notes(s, "Collection is worthless if it sits in a database. Every day the system surfaces one story, ready to post / testify.")

# ============ SLIDE 10 — IMPACT & COST ============
s = slide(INK)
header_bar(s, "Impact & cost", "Built for a nonprofit budget", dark=True)
stats = [("3,700+", "stories"), ("120+", "countries"), ("20", "languages"), ("~2.1¢", "per story")]
cw = 2.7
x0 = (SW - (cw * 4 + 0.35 * 3)) / 2
for i, (n, lab) in enumerate(stats):
    x = x0 + i * (cw + 0.35)
    rect(s, x, 1.5, cw, 1.9, fill=NAVY, line=RGBColor(0x2A, 0x3A, 0x7A))
    text(s, x, 1.65, cw, 1.1, [[H(n, 40, EMER)]], align=PP_ALIGN.CENTER)
    text(s, x, 2.8, cw, 0.5, [[B(lab, 14, LBLUE)]], align=PP_ALIGN.CENTER)
picture(s, IMG.format("cost"), 0.45, 3.7, 8.4, 3.35)
text(s, 9.2, 3.8, 3.8, 3.2, [
    [H("≈ $79", 34, GOLD)],
    [B("to process the entire archive.", 14, LBLUE)],
    [B("", 8, WHITE)],
    [B("$150,000 = one NY Proud campaign. We're a rounding error of that.", 14, WHITE, True)],
], anchor=MSO_ANCHOR.MIDDLE, space=1.08)
notes(s, "At NYIC's own '40 Years, 40,000 Stories' vision, the whole AI pipeline costs under $850 — half a percent of one "
         "billboard campaign. One part-time coordinator + volunteers does the rest.")

# ============ SLIDE 11 — WHY IT WINS ============
s = slide(MIST)
header_bar(s, "Deployable Monday — and built to last")
pts = [
    ("Demo-proof", "Runs offline-first on seeded data — no live API can break the pitch."),
    ("No vendor lock", "Supabase + open-source Leaflet/OpenStreetMap maps."),
    ("Volunteer-powered", "Code-4-Good & campus partners staff the recording booths."),
    ("Coalition-native", "Every one of 200+ member orgs is credited and sees its own stories."),
]
for i, (t, d) in enumerate(pts):
    x = 1.4 + (i % 2) * 5.7
    y = 1.9 + (i // 2) * 1.9
    rect(s, x, y, 5.2, 1.55, fill=WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    rect(s, x, y, 0.14, 1.55, fill=[NAVY, EMER, GOLD, RED][i])
    text(s, x + 0.4, y + 0.22, 4.6, 1.2, [
        [H(t, 17, NAVY)],
        [B(d, 13.5, SLATE)],
    ], space=1.05)
text(s, 0, 5.9, SW, 0.8, [[B("4.5M immigrants served  ·  40 years  ·  a coalition of 200+ organizations", 15, SLATE, True)]], align=PP_ALIGN.CENTER)
notes(s, "Sustainability + partnership story. This isn't a hack-week toy — it's deployable Monday.")

# ============ SLIDE 12 — CLOSE ============
s = slide(NAVY)
picture(s, IMG.format("origins"), 0, 3.6, SW, 3.9)
rect(s, 0, 3.6, SW, 3.9, fill=NAVY)
text(s, 1.0, 1.4, 11.3, 2.4, [
    [H("“We don't ask immigrants to prove they belong.", 32, WHITE)],
    [H("We ask them to choose us back.”", 32, GOLD)],
], space=1.1)
text(s, 1.0, 4.3, 11, 0.7, [[H("#IChooseNY   ·   #IChooseYou", 22, WHITE)]])
text(s, 1.0, 5.2, 11, 0.6, [[B("ichooseny.netlify.app", 18, EMER, True)]])
rect(s, 10.8, 5.0, 1.7, 1.7, fill=WHITE)
text(s, 10.8, 5.0, 1.7, 1.7, [[B("[ QR code\nto site ]", 11, SLATE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
notes(s, "The memorized close. End on the tagline, then SILENCE. Don't keep talking.")

out = "/Users/hasanmalik/Desktop/nyic/ichooseny-deck.pptx"
prs.save(out)
print("saved", out, "—", len(prs.slides._sldIdLst), "slides")
